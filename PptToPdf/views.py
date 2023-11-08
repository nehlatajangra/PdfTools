from django.shortcuts import render
from reportlab.pdfgen import canvas
from django.shortcuts import render
from pptx import Presentation
from reportlab.pdfgen import canvas
from io import BytesIO
from django.http import HttpResponse, HttpResponseBadRequest
# from aspose.cells import Workbook


def pptToPdfView(request):
    if request.method == 'POST' and request.FILES['ppt_file']:
        uploaded_file = request.FILES['ppt_file']

        # Check if the file is a PPT file
        if not uploaded_file.name.lower().endswith('.pptx'):
            return HttpResponseBadRequest("Invalid file format. Only PPTX files are supported.")

        # Create a PowerPoint presentation object
        ppt_buffer = BytesIO(uploaded_file.read())
        presentation = Presentation(ppt_buffer)

        # Create a PDF file in memory using ReportLab
        pdf_buffer = BytesIO()
        pdf = canvas.Canvas(pdf_buffer)

        # Iterate through slides and add content to the PDF
        for slide in presentation.slides:
            print(slide)
            for shape in slide.shapes:
                print(shape)
                print(shape.has_text_frame)
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        print(paragraph)
                        for run in paragraph.runs:
                            print(run)
                            pdf.drawString(100, 700, run.text)
            pdf.showPage()

        pdf.save()

        # Create the HTTP response with the PDF content
        pdf_buffer.seek(0)
        response = HttpResponse(pdf_buffer.read(), content_type='application/pdf')
        response['Content-Disposition'] = 'attachment; filename=converted_ppt.pdf'

        return response

    return render(request, 'pptToPdf.html')